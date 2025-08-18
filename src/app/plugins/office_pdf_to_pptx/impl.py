# src/app/plugins/office_pdf_to_pptx/impl.py
from __future__ import annotations

import os
import io
import uuid
import shutil
import tempfile
from typing import Any, Dict, List, Tuple

# deps: httpx, pillow, pymupdf (fitz), python-pptx
import httpx
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt

try:
    import fitz  # PyMuPDF
except Exception as e:  # pragma: no cover
    fitz = None

try:
    # If your project already has helpers, you can uncomment and reuse them:
    # from app.tools.office_io import save_artifact as _save_artifact_builtin
    HAVE_OFFICE_IO = False
except Exception:
    HAVE_OFFICE_IO = False

from app.core.config import settings


# ---------------------------
# Utilities (self-contained)
# ---------------------------

def _emit(ctx, step: str, status: str = "info", payload: Dict[str, Any] | None = None):
    """Best-effort event emission; safe if events bridge changes."""
    try:
        ctx.events.emit(ctx.tenant_id, ctx.job_id, step, status, payload or {})
    except Exception:
        try:
            # Older signature: emit(job_id, step, payload)
            ctx.events.emit(ctx.job_id, step, payload or {})
        except Exception:
            pass


def _download_to_tmp(url: str, tmpdir: str) -> str:
    filename = f"src_{uuid.uuid4().hex}.pdf"
    dest = os.path.join(tmpdir, filename)
    with httpx.stream("GET", url, follow_redirects=True, timeout=60) as r:
        r.raise_for_status()
        with open(dest, "wb") as f:
            for chunk in r.iter_bytes():
                f.write(chunk)
    return dest


def _pdf_to_images(pdf_path: str, tmpdir: str, dpi: int, max_slides: int, img_fmt: str) -> List[str]:
    """
    Render PDF pages to images. Prefers PyMuPDF. If unavailable, raises a clear error.
    """
    if fitz is None:
        raise RuntimeError(
            "PyMuPDF (fitz) is required for PDF rendering. "
            "Please add 'pymupdf' to your requirements."
        )

    doc = fitz.open(pdf_path)
    img_paths: List[str] = []
    count = min(len(doc), max_slides if max_slides else len(doc))

    for i in range(count):
        page = doc.load_page(i)
        # DPI -> zoom
        zoom = dpi / 72.0
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat, alpha=False)
        img_bytes = pix.tobytes("png") if img_fmt.upper() == "PNG" else pix.tobytes("jpeg")
        ext = "png" if img_fmt.upper() == "PNG" else "jpg"
        out_path = os.path.join(tmpdir, f"page_{i+1:04d}.{ext}")
        with open(out_path, "wb") as f:
            f.write(img_bytes)
        img_paths.append(out_path)

    doc.close()
    return img_paths


def _images_to_pptx(
    image_paths: List[str],
    title: str,
    fit_mode: str = "contain",
) -> Presentation:
    """
    Build a simple deck: 1 title slide (if any) + 1 image per slide, full-bleed.
    """
    prs = Presentation()

    # Title slide
    if title:
        layout = prs.slide_layouts[0]  # title + subtitle
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = title
        if s.placeholders and len(s.placeholders) > 1:
            s.placeholders[1].text = ""

    # Use slide size similar to 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for p in image_paths:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        with Image.open(p) as im:
            img_w_px, img_h_px = im.size
        slide_w_in = prs.slide_width.inches
        slide_h_in = prs.slide_height.inches

        # Convert pixels to inches assuming 96 dpi (typical screen)
        img_w_in = img_w_px / 96.0
        img_h_in = img_h_px / 96.0

        # Compute scale to fit or cover
        scale_w = slide_w_in / img_w_in
        scale_h = slide_h_in / img_h_in
        if fit_mode == "cover":
            scale = max(scale_w, scale_h)
        else:
            scale = min(scale_w, scale_h)

        w = Inches(img_w_in * scale)
        h = Inches(img_h_in * scale)
        x = (prs.slide_width - w) / 2
        y = (prs.slide_height - h) / 2

        slide.shapes.add_picture(p, x, y, width=w, height=h)

    return prs


def _save_artifact_local(local_path: str, preferred_name: str) -> Tuple[str, str]:
    """
    Save under artifacts dir and build a public URL.
    Uses env or settings:
      - ARTIFACTS_DIR      (default: /data/artifacts)
      - PUBLIC_BASE_URL    (default: http://localhost:8080)
      - ARTIFACTS_HTTP_BASE (optional: override '/artifacts')
    """
    artifacts_dir = os.getenv("ARTIFACTS_DIR", "/data/artifacts")
    http_base = os.getenv("ARTIFACTS_HTTP_BASE", "/artifacts")
    base_url = getattr(settings, "PUBLIC_BASE_URL", None) or os.getenv("PUBLIC_BASE_URL", "http://localhost:8080")

    os.makedirs(artifacts_dir, exist_ok=True)
    # Unique filename
    name = f"{os.path.splitext(preferred_name)[0]}_{uuid.uuid4().hex}.pptx"
    dest = os.path.join(artifacts_dir, name)
    shutil.copy2(local_path, dest)

    url = f"{base_url.rstrip('/')}{http_base}{'/' if not http_base.endswith('/') else ''}{name}"
    return dest, url


# ---------------------------
# Entry point
# ---------------------------

def run(inputs: Dict[str, Any], ctx) -> Dict[str, Any]:
    """
    Contract:
      - inputs validated by manifest JSON-Schema
      - ctx has: tenant_id, job_id, events
      - returns: { pptx_url, pptx_key }
    """
    file_url: str = inputs["file_url"]
    title: str = inputs["title"]
    max_slides: int = int(inputs.get("max_slides", 50))
    dpi: int = int(inputs.get("dpi", 140))
    fit_mode: str = inputs.get("fit_mode", "contain")
    image_format: str = inputs.get("image_format", "PNG").upper()

    _emit(ctx, "pdf.fetch", "started", {"url": file_url})

    with tempfile.TemporaryDirectory(prefix="pdf2pptx_") as tmp:
        # 1) Download PDF
        pdf_path = _download_to_tmp(file_url, tmp)
        _emit(ctx, "pdf.fetch", "succeeded", {"path": pdf_path})

        # 2) Render pages to images
        _emit(ctx, "pdf.render", "started", {"dpi": dpi, "max_slides": max_slides, "format": image_format})
        pages = _pdf_to_images(pdf_path, tmp, dpi=dpi, max_slides=max_slides, img_fmt=image_format)
        _emit(ctx, "pdf.render", "succeeded", {"count": len(pages)})

        if not pages:
            raise RuntimeError("No pages found in PDF")

        # 3) Build PPTX
        _emit(ctx, "pptx.build", "started", {"slides": len(pages)})
        prs = _images_to_pptx(pages, title=title, fit_mode=fit_mode)

        out_path = os.path.join(tmp, "slides.pptx")
        prs.save(out_path)
        _emit(ctx, "pptx.build", "succeeded", {"path": out_path})

        # 4) Save artifact
        pptx_key, pptx_url = _save_artifact_local(out_path, "slides.pptx")
        # If you already have a builtin artifact saver, switch to:
        # if HAVE_OFFICE_IO:
        #     key, url = _save_artifact_builtin(out_path, "slides.pptx")
        # else:
        #     key, url = _save_artifact_local(out_path, "slides.pptx")

    return {"pptx_key": pptx_key, "pptx_url": pptx_url}
