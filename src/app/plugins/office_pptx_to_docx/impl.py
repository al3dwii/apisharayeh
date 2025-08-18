# src/app/plugins/office_pptx_to_docx/impl.py
from __future__ import annotations

import os
import io
import uuid
import shutil
import tempfile
from typing import Any, Dict, List, Tuple

import httpx
from app.core.config import settings

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from PIL import Image  # used only if include_images=True


# ---------------------------
# Event + artifact helpers
# ---------------------------

def _emit(ctx, step: str, status: str = "info", payload: Dict[str, Any] | None = None):
    """Best-effort event emission; compatible with multiple bridge shapes."""
    try:
        ctx.events.emit(ctx.tenant_id, ctx.job_id, step, status, payload or {})
    except Exception:
        try:
            ctx.events.emit(ctx.job_id, step, payload or {})
        except Exception:
            pass


def _save_artifact_local(local_path: str, preferred_name: str) -> Tuple[str, str]:
    """
    Save under artifacts dir and build a public URL.
    Env/settings:
      - ARTIFACTS_DIR (default /data/artifacts)
      - PUBLIC_BASE_URL (http://localhost:8080)
      - ARTIFACTS_HTTP_BASE (/artifacts)
    """
    artifacts_dir = os.getenv("ARTIFACTS_DIR", "/data/artifacts")
    http_base = os.getenv("ARTIFACTS_HTTP_BASE", "/artifacts")
    base_url = getattr(settings, "PUBLIC_BASE_URL", None) or os.getenv("PUBLIC_BASE_URL", "http://localhost:8080")

    os.makedirs(artifacts_dir, exist_ok=True)
    name = f"{os.path.splitext(preferred_name)[0]}_{uuid.uuid4().hex}.docx"
    dest = os.path.join(artifacts_dir, name)
    shutil.copy2(local_path, dest)

    sep = "" if http_base.endswith("/") else "/"
    url = f"{base_url.rstrip('/')}{http_base}{sep}{name}"
    return dest, url


def _download_pptx(url: str, tmpdir: str) -> str:
    name = f"src_{uuid.uuid4().hex}.pptx"
    dest = os.path.join(tmpdir, name)
    with httpx.stream("GET", url, follow_redirects=True, timeout=60) as r:
        r.raise_for_status()
        with open(dest, "wb") as f:
            for chunk in r.iter_bytes():
                f.write(chunk)
    return dest


# ---------------------------
# PPTX -> outline extraction
# ---------------------------

def _shape_texts(shape) -> List[Tuple[int, str]]:
    """
    Extract (level, text) tuples from a text-bearing shape.
    For tables, we flatten cells row-by-row at level 0.
    """
    items: List[Tuple[int, str]] = []
    # Text frames
    if hasattr(shape, "has_text_frame") and shape.has_text_frame and shape.text_frame:
        tf = shape.text_frame
        for p in tf.paragraphs:
            # Gather runs (maintain soft formatting losslessly as plain text)
            text = "".join(run.text for run in p.runs).strip()
            if not text:
                continue
            level = getattr(p, "level", 0) or 0
            items.append((int(level), text))
    # Tables
    elif hasattr(shape, "has_table") and shape.has_table and shape.table:
        table = shape.table
        for row in table.rows:
            row_text = []
            for cell in row.cells:
                t = (cell.text or "").strip()
                if t:
                    row_text.append(t)
            if row_text:
                items.append((0, " | ".join(row_text)))
    return items


def _slide_heading(slide, fallback_prefix: str, idx: int) -> str:
    """
    Prefer the slide title placeholder. Otherwise, use first non-empty paragraph
    as heading. Fallback to 'Slide N'.
    """
    try:
        if slide.shapes.title and slide.shapes.title.text:
            t = slide.shapes.title.text.strip()
            if t:
                return t
    except Exception:
        pass

    # Find first non-empty text in any shape
    for shp in slide.shapes:
        for _, txt in _shape_texts(shp):
            if txt:
                return txt[:120]  # trim long headings

    return f"{fallback_prefix} {idx}"


def _collect_images(slide, tmpdir: str) -> List[str]:
    """
    Extract slide images to tmp files; return their paths.
    """
    out: List[str] = []
    for shp in slide.shapes:
        if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                image = shp.image  # python-pptx image object
                ext = image.ext or "png"
                name = f"img_{uuid.uuid4().hex}.{ext}"
                path = os.path.join(tmpdir, name)
                with open(path, "wb") as f:
                    f.write(image.blob)
                out.append(path)
            except Exception:
                continue
    return out


# ---------------------------
# DOCX building
# ---------------------------

def _apply_doc_defaults(doc: Document) -> None:
    # Reasonable defaults; keep it simple and readable.
    style = doc.styles["Normal"]
    font = style.font
    font.name = "Calibri"
    font.size = Pt(11)


def _add_bullets(doc: Document, items: List[Tuple[int, str]]) -> None:
    """
    Add bullet paragraphs, using 'List Bullet' and left indentation per level.
    python-docx doesn't support true multilevel numbering styles declaratively;
    we emulate with indentation.
    """
    for level, text in items:
        p = doc.add_paragraph(text)
        p.style = "List Bullet"
        # 0.25" indent per level (tweak as needed)
        indent_in = max(0, min(6, level)) * 0.25
        p.paragraph_format.left_indent = Inches(indent_in)


def _add_images(doc: Document, image_paths: List[str]) -> None:
    """
    Insert images scaled to fit ~6.5" width while keeping aspect ratio.
    """
    max_w_in = 6.5
    for path in image_paths:
        try:
            with Image.open(path) as im:
                w, h = im.size
                # Rough DPI assumption for sizing; python-docx scales by width
                doc.add_picture(path, width=Inches(max_w_in))
        except Exception:
            continue


# ---------------------------
# Entry point
# ---------------------------

def run(inputs: Dict[str, Any], ctx) -> Dict[str, Any]:
    """
    Inputs:
      - file_url (str)            : PPTX URL
      - include_notes (bool)      : default False
      - include_images (bool)     : default False
      - max_slides (int)          : default 500
      - heading_prefix (str)      : default 'Slide'
    Returns:
      - { docx_key, docx_url }
    """
    file_url: str = inputs["file_url"]
    include_notes: bool = bool(inputs.get("include_notes", False))
    include_images: bool = bool(inputs.get("include_images", False))
    max_slides: int = int(inputs.get("max_slides", 500))
    heading_prefix: str = str(inputs.get("heading_prefix", "Slide"))

    _emit(ctx, "pptx.fetch", "started", {"url": file_url})

    with tempfile.TemporaryDirectory(prefix="pptx2docx_") as tmp:
        # 1) Download PPTX
        pptx_path = _download_pptx(file_url, tmp)
        _emit(ctx, "pptx.fetch", "succeeded", {"path": pptx_path})

        # 2) Load presentation
        prs = Presentation(pptx_path)
        slide_count = min(len(prs.slides), max_slides)
        _emit(ctx, "pptx.parse", "started", {"slides": slide_count})

        # 3) Build DOCX
        doc = Document()
        _apply_doc_defaults(doc)

        # Optional: add a document title from first slide
        try:
            if slide_count > 0:
                title_guess = _slide_heading(prs.slides[0], heading_prefix, 1)
                h = doc.add_heading(title_guess, level=0)  # title style
                h.alignment = WD_ALIGN_PARAGRAPH.CENTER
        except Exception:
            pass

        for i in range(slide_count):
            slide = prs.slides[i]
            # Heading
            heading = _slide_heading(slide, heading_prefix, i + 1)
            doc.add_heading(f"{heading_prefix} {i+1}: {heading}", level=1)

            # Images (optional)
            if include_images:
                imgs = _collect_images(slide, tmp)
                if imgs:
                    _add_images(doc, imgs)

            # Text bullets
            bullets: List[Tuple[int, str]] = []
            for shp in slide.shapes:
                # Skip the title placeholder to avoid duplication
                try:
                    if shp == slide.shapes.title:
                        continue
                except Exception:
                    pass
                bullets.extend(_shape_texts(shp))
            if bullets:
                _add_bullets(doc, bullets)

            # Speaker notes (optional)
            if include_notes:
                try:
                    if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                        notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
                        if notes_text:
                            doc.add_heading("Notes", level=2)
                            for line in notes_text.splitlines():
                                line = line.strip()
                                if line:
                                    p = doc.add_paragraph(line)
                                    p.style = "List Paragraph"
                except Exception:
                    pass

            # Spacer between slides
            doc.add_paragraph("")

        # 4) Save DOCX
        out_path = os.path.join(tmp, "presentation.docx")
        doc.save(out_path)
        _emit(ctx, "docx.build", "succeeded", {"path": out_path})

        # 5) Save artifact
        docx_key, docx_url = _save_artifact_local(out_path, "presentation.docx")

    return {"docx_key": docx_key, "docx_url": docx_url}
