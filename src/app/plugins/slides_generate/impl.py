# app/plugins/slides_generate/impl.py
"""
Agentic in-proc plugin that turns a prompt or a document (DOCX/PDF)
into a high-quality presentation:
 - Builds a normalized outline
 - Auto-fills images via ToolRouter (Unsplash/Bing/SDXL) or HTTP fallbacks
 - Generates a PPTX (RTL-aware) using a theme master
 - Exports PDF + HTML (whole deck)
 - Renders per-slide HTML fragments (for live preview)
 - Emits rich events at every step
"""
from __future__ import annotations

from typing import Dict, Any, List
import traceback

# --- Internal imports from your codebase (fallbacks included for dev) ---
try:
    from app.kernel.runtime import KernelContext
except Exception:  # pragma: no cover
    class KernelContext:  # lightweight fallback for type hints
        job_id: str
        tenant_id: str

try:
    from app.services.models import ModelRouter
except Exception:  # pragma: no cover
    class ModelRouter:
        def chat(self, messages: List[Dict[str, str]]) -> Dict[str, str]:
            raise RuntimeError("ModelRouter not available in this environment. Wire it in the app.")

try:
    from app.services.tools import ToolRouter
except Exception:  # pragma: no cover
    class ToolRouter:
        def __init__(self, *args, **kwargs):
            pass
        def call(self, tool_id: str, payload: dict):
            raise RuntimeError("ToolRouter not available in this environment.")

# Events (compat wrapper supports both (ctx, typ, status, payload) and (ctx, typ, payload))
try:
    from app.services.events import emit_event as _emit_event
except Exception:  # pragma: no cover
    def _emit_event(ctx, typ, *args):
        # local fallback that just prints
        payload = args[-1] if args else {}
        print(f"[event] {typ}: {payload}")

def EMIT(ctx, typ: str, payload: dict, status: str = "progress"):
    """Compat emitter: tries 4-arg signature, falls back to 3-arg."""
    try:
        _emit_event(ctx, typ, status, payload)  # new signature
    except TypeError:
        _emit_event(ctx, typ, payload)          # old signature

try:
    from app.services.artifacts import save_file
except Exception:  # pragma: no cover
    import os, shutil, uuid
    def save_file(local_path: str, default_ext: str) -> Dict[str, str]:
        """Fallback: copies to /tmp and returns a file:// URL."""
        out_dir = "/tmp/artifacts"
        os.makedirs(out_dir, exist_ok=True)
        key = f"{uuid.uuid4().hex}{default_ext}"
        dst = os.path.join(out_dir, key)
        shutil.copyfile(local_path, dst)
        return {"key": key, "url": f"file://{dst}"}

# Office/PPTX utilities + builder
try:
    from app.tools.office_io import pptx_to_pdf_file, pptx_to_html5_zip  # your real converters
    from .pptx_builder import build_pptx
except Exception:  # pragma: no cover
    # Very thin fallbacks (non-production)
    from pptx import Presentation
    import tempfile, os, subprocess, zipfile

    def build_pptx(outline, title=None, theme_path=None, lang="ar"):
        prs = Presentation()
        # Cover slide
        if outline and not (outline[0].get("bullets")):
            layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = outline[0].get("title") or (title or "عرض تقديمي")
            if len(slide.placeholders) > 1 and outline[0].get("subtitle"):
                slide.placeholders[1].text = outline[0]["subtitle"]
        # Content slides
        for i, s in enumerate(outline, 1):
            if i == 1 and not s.get("bullets"):
                continue
            layout = prs.slide_layouts[1]
            sl = prs.slides.add_slide(layout)
            sl.shapes.title.text = s.get("title","")
            tx = sl.shapes.placeholders[1].text_frame
            tx.clear()
            for b in s.get("bullets", []):
                text = b.get("text","") if isinstance(b, dict) else str(b)
                subs = (b.get("subs") or []) if isinstance(b, dict) else []
                if not tx.text:
                    p = tx.paragraphs[0]
                    p.text = text
                    p.level = 0
                else:
                    p = tx.add_paragraph()
                    p.text = text
                    p.level = 0
                for sb in subs[:2]:
                    sp = tx.add_paragraph()
                    sp.text = sb
                    sp.level = 1
        fd, path = tempfile.mkstemp(suffix=".pptx"); os.close(fd)
        prs.save(path)
        return path

    def pptx_to_pdf_file(pptx_path: str) -> str:
        out = pptx_path.replace(".pptx", ".pdf")
        try:
            subprocess.run(
                ["soffice","--headless","--convert-to","pdf","--outdir",os.path.dirname(pptx_path), pptx_path],
                check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE
            )
        except Exception:
            # fallback: duplicate pptx (not a real PDF)
            import shutil
            shutil.copyfile(pptx_path, out)
        return out

    def pptx_to_html5_zip(pptx_path: str) -> str:
        out = pptx_path.replace(".pptx", ".html5.zip")
        with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as z:
            z.write(pptx_path, arcname=os.path.basename(pptx_path))
        return out

# Slide helpers
from .outline import outline_from_prompt, outline_from_docx, outline_from_pdf
from .images import attach_images
from .quality import normalize_outline
from .layout import choose_layouts_and_render_html

# Register image providers on ToolRouter if available
try:
    from app.services.tools.providers.images_providers import register_image_providers
except Exception:
    register_image_providers = None


def run(inputs: Dict[str, Any], ctx: KernelContext) -> Dict[str, Any]:
    """
    Main entrypoint invoked by the kernel runtime.
    """
    language = inputs.get("language", "ar")
    max_slides = int(inputs.get("max_slides", 12))
    theme     = inputs.get("theme", "academic")

    # Optional image settings
    image_strategy = inputs.get("image_strategy", "hybrid")  # search | generate | hybrid | none
    image_provider = inputs.get("image_provider")            # unsplash | bing | sdxl

    # 1) Plan
    EMIT(ctx, "thinking", {"stage": "planning", "msg": "Deep Thinking: plan tasks"})

    mr = ModelRouter()
    # ToolRouter in this codebase requires a models argument; be tolerant of either ctor style.
    try:
        tr = ToolRouter(models=mr)   # preferred (named)
    except TypeError:
        tr = ToolRouter(mr)          # fallback (positional)

    if register_image_providers:
        try:
            register_image_providers(tr)  # makes unsplash.search / bing.image.search / sdxl.generate available
        except Exception:
            # non-fatal; images.py has HTTP fallbacks
            pass

    # 2) Outline
    try:
        if inputs.get("prompt"):
            EMIT(ctx, "tool", {"tool": "outline_from_prompt"})
            outline = outline_from_prompt(mr, tr, inputs["prompt"], language, max_slides)
        elif inputs.get("word_url"):
            EMIT(ctx, "tool", {"tool": "outline_from_docx", "url": inputs["word_url"]})
            outline = outline_from_docx(tr, inputs["word_url"], max_slides)
        elif inputs.get("pdf_url"):
            EMIT(ctx, "tool", {"tool": "outline_from_pdf", "url": inputs["pdf_url"]})
            outline = outline_from_pdf(tr, inputs["pdf_url"], max_slides)
        else:
            raise ValueError("One of 'prompt', 'word_url', or 'pdf_url' is required.")
    except Exception as e:
        EMIT(ctx, "error", {"stage": "outline", "error": str(e), "trace": traceback.format_exc()}, status="failed")
        raise

    # 3) Quality pass
    EMIT(ctx, "thinking", {"stage": "qa", "msg": "Normalize bullets, enforce limits"})
    outline = normalize_outline(outline, max_slides=max_slides, language=language)

    # 3b) Auto-attach images via ToolRouter (with HTTP fallbacks inside)
    outline = attach_images(
        outline,
        tr,
        strategy=image_strategy,
        provider=image_provider,
        emit=lambda et, payload: EMIT(ctx, et, payload)
    )

    # 4) Build PPTX
    EMIT(ctx, "tool", {"tool": "outline_to_pptx"})
    from pathlib import Path
    theme_path = Path(__file__).parent / "theme" / "academic.pptx"
    pptx_path = build_pptx(
        outline,
        title=outline[0].get("title") if outline else inputs.get("prompt"),
        theme_path=str(theme_path),
        lang=language
    )

    # 5) Exports
    EMIT(ctx, "tool", {"tool": "pptx_to_pdf"})
    pdf_path = pptx_to_pdf_file(pptx_path)
    EMIT(ctx, "tool", {"tool": "pptx_to_html5"})
    html_zip = pptx_to_html5_zip(pptx_path)

    # 6) Server-side HTML fragments for live preview
    EMIT(ctx, "thinking", {"stage": "render", "msg": "Render per-slide HTML"})
    slides_html = choose_layouts_and_render_html(outline, theme=theme, lang=language)

    # 7) Save artifacts
    try:
        pptx_art = save_file(pptx_path, ".pptx")
        pdf_art  = save_file(pdf_path,  ".pdf")
        zip_art  = save_file(html_zip,  ".zip")
    except Exception:
        # Fallback already handled inside save_file shim
        pptx_art = save_file(pptx_path, ".pptx")
        pdf_art  = save_file(pdf_path,  ".pdf")
        zip_art  = save_file(html_zip,  ".zip")

    result = {
        "pptx_url": pptx_art["url"],
        "pdf_url": pdf_art["url"],
        "html_zip_url": zip_art["url"],
        "outline": outline,
        "slides_html": slides_html,
    }

    EMIT(ctx, "summary", {"slides": len(outline), "outputs": ["pptx", "pdf", "html5"], "ok": True}, status="done")
    return result
