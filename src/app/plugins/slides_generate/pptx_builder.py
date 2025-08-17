
# app/plugins/slides_generate/pptx_builder.py
from __future__ import annotations
from typing import List, Dict, Any, Optional
import os, io, tempfile, requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

DEFAULT_FONT = "Tajawal"

def _set_run_style(run, font_size: int = 24, bold: bool = False, color: Optional[RGBColor] = None, font_name: str = DEFAULT_FONT):
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color

def _add_bullet_paragraph(tf, text: str, level: int, rtl: bool, font_size: int = 24):
    if len(tf.paragraphs) == 1 and not tf.paragraphs[0].text:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.level = level
    p.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT if rtl else PP_PARAGRAPH_ALIGNMENT.LEFT
    for r in p.runs:
        _set_run_style(r, font_size=font_size)

def _font_size_for_bullet_count(count: int) -> int:
    if count <= 3: return 28
    if count == 4: return 26
    if count == 5: return 24
    return 22

def _download_image(url: str) -> Optional[str]:
    try:
        r = requests.get(url, timeout=10)
        r.raise_for_status()
        fd, path = tempfile.mkstemp(suffix=".img")
        os.write(fd, r.content)
        os.close(fd)
        return path
    except Exception:
        return None

def _title_slide(prs: Presentation, title: str, subtitle: Optional[str], rtl: bool):
    layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if subtitle is not None and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
    # apply font and alignment
    t = slide.shapes.title.text_frame.paragraphs[0]
    t.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT if rtl else PP_PARAGRAPH_ALIGNMENT.LEFT
    for r in slide.shapes.title.text_frame.paragraphs[0].runs:
        _set_run_style(r, font_size=44, bold=True)
    if subtitle is not None and len(slide.placeholders) > 1:
        p = slide.placeholders[1].text_frame.paragraphs[0]
        p.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT if rtl else PP_PARAGRAPH_ALIGNMENT.LEFT
        for r in p.runs:
            _set_run_style(r, font_size=22)

def _list_slide(prs: Presentation, slide_data: Dict[str, Any]):
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    title_shape = slide.shapes.title
    body = slide.shapes.placeholders[1]
    rtl = slide_data.get("rtl", False)
    title_shape.text = slide_data.get("title","")
    title_p = title_shape.text_frame.paragraphs[0]
    title_p.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT if rtl else PP_PARAGRAPH_ALIGNMENT.LEFT
    for r in title_p.runs:
        _set_run_style(r, font_size=36, bold=True)

    tf = body.text_frame
    tf.clear()
    bullets = slide_data.get("bullets", [])
    size = _font_size_for_bullet_count(len(bullets))
    for b in bullets:
        text = b["text"] if isinstance(b, dict) else str(b)
        _add_bullet_paragraph(tf, text, level=0, rtl=rtl, font_size=size)
        subs = b.get("subs", []) if isinstance(b, dict) else []
        for sb in subs[:2]:
            _add_bullet_paragraph(tf, sb, level=1, rtl=rtl, font_size=max(size-2, 18))

def _two_col_slide(prs: Presentation, slide_data: Dict[str, Any]):
    layout = prs.slide_layouts[3]  # Two Content
    slide = prs.slides.add_slide(layout)
    rtl = slide_data.get("rtl", False)
    slide.shapes.title.text = slide_data.get("title","")
    title_p = slide.shapes.title.text_frame.paragraphs[0]
    title_p.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT if rtl else PP_PARAGRAPH_ALIGNMENT.LEFT
    for r in title_p.runs:
        _set_run_style(r, font_size=34, bold=True)

    left_pl = slide.placeholders[1]
    right_pl = slide.placeholders[2]

    bullets = slide_data.get("bullets", [])
    mid = (len(bullets)+1)//2
    left, right = bullets[:mid], bullets[mid:]

    ltf = left_pl.text_frame; ltf.clear()
    rtf = right_pl.text_frame; rtf.clear()
    size = _font_size_for_bullet_count(len(bullets))

    for b in left:
        text = b["text"] if isinstance(b, dict) else str(b)
        _add_bullet_paragraph(ltf, text, level=0, rtl=rtl, font_size=size)
        for sb in (b.get("subs", []) if isinstance(b, dict) else [])[:2]:
            _add_bullet_paragraph(ltf, sb, level=1, rtl=rtl, font_size=max(size-2, 18))

    for b in right:
        text = b["text"] if isinstance(b, dict) else str(b)
        _add_bullet_paragraph(rtf, text, level=0, rtl=rtl, font_size=size)
        for sb in (b.get("subs", []) if isinstance(b, dict) else [])[:2]:
            _add_bullet_paragraph(rtf, sb, level=1, rtl=rtl, font_size=max(size-2, 18))

def _title_image_slide(prs: Presentation, slide_data: Dict[str, Any]):
    # Use Title Only layout then draw picture + textbox
    layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    rtl = slide_data.get("rtl", False)
    # Title at top
    if slide.shapes.title:
        slide.shapes.title.text = slide_data.get("title","")
        p = slide.shapes.title.text_frame.paragraphs[0]
        p.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT if rtl else PP_PARAGRAPH_ALIGNMENT.LEFT
        for r in p.runs:
            _set_run_style(r, font_size=34, bold=True)

    # Image on right/left depending on RTL
    img_url = slide_data.get("image")
    img_path = _download_image(img_url) if isinstance(img_url, str) else None
    img_w = Inches(5.3)
    img_h = Inches(4.0)
    img_left = Inches(0.6) if rtl else Inches(6.0 - 0.6 - 0.1)  # approximate
    img_top = Inches(2.0)

    text_left = Inches(6.4) if rtl else Inches(0.6)
    text_top = Inches(2.0)
    text_w = Inches(5.2)
    text_h = Inches(3.9)

    if img_path and os.path.exists(img_path):
        slide.shapes.add_picture(img_path, img_left, img_top, width=img_w, height=img_h)
    # Text box for bullets
    tb = slide.shapes.add_textbox(text_left, text_top, text_w, text_h)
    tf = tb.text_frame
    tf.clear()
    bullets = slide_data.get("bullets", [])
    size = _font_size_for_bullet_count(len(bullets))
    for b in bullets:
        text = b["text"] if isinstance(b, dict) else str(b)
        _add_bullet_paragraph(tf, text, level=0, rtl=rtl, font_size=size)
        subs = b.get("subs", []) if isinstance(b, dict) else []
        for sb in subs[:2]:
            _add_bullet_paragraph(tf, sb, level=1, rtl=rtl, font_size=max(size-2, 18))

def build_pptx(outline: List[Dict[str, Any]], title: Optional[str] = None, theme_path: Optional[str] = None, lang: str = "ar") -> str:
    """
    Build a PPTX with decent defaults, RTL-aware, dynamic font sizes and simple layouts.
    """
    prs = Presentation(theme_path) if theme_path and os.path.exists(theme_path) else Presentation()

    # cover
    if outline:
        s0 = outline[0]
        if not s0.get("bullets"):
            _title_slide(prs, s0.get("title") or title or "عرض تقديمي", s0.get("subtitle"), rtl=s0.get("rtl", False))

    # rest
    for idx, s in enumerate(outline, 1):
        if idx == 1 and not s.get("bullets"):
            continue  # already added cover
        # choose layout
        layout_hint = s.get("layout_hint","auto")
        bullets = s.get("bullets") or []
        if layout_hint == "title_image" or s.get("image"):
            _title_image_slide(prs, s)
        elif layout_hint == "two_col" or len(bullets) > 4:
            _two_col_slide(prs, s)
        elif layout_hint == "title" and not bullets:
            _title_slide(prs, s.get("title") or "", s.get("subtitle"), rtl=s.get("rtl", False))
        else:
            _list_slide(prs, s)

    fd, out_path = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)
    prs.save(out_path)
    return out_path
