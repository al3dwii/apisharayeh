import asyncio
import io
import json
import uuid
import os
from typing import Optional

import httpx
from pdfminer.high_level import extract_text
from pptx import Presentation
from pptx.util import Pt

from app.core.config import settings

try:
    import aioboto3
except Exception:
    aioboto3 = None  # S3 optional


# ---------- helpers ----------

async def _fetch_bytes_from_url(url: str, timeout: int = 90) -> bytes:
    async with httpx.AsyncClient(timeout=timeout) as c:
        r = await c.get(url)
        r.raise_for_status()
        return r.content


async def _fetch_bytes_from_s3(key: str) -> bytes:
    if not (settings.S3_BUCKET and aioboto3):
        raise RuntimeError("S3 not configured")
    session = aioboto3.Session()
    async with session.client("s3", region_name=settings.AWS_REGION) as s3:
        obj = await s3.get_object(Bucket=settings.S3_BUCKET, Key=key)
        return await obj["Body"].read()


async def _extract_text_async(pdf_bytes: bytes) -> str:
    def _run(buf: bytes) -> str:
        return extract_text(io.BytesIO(buf)) or ""
    return await asyncio.to_thread(_run, pdf_bytes)



async def _upload_bytes(
    data: bytes,
    key: str,
    content_type: str = "application/vnd.openxmlformats-officedocument.presentationml.presentation",
) -> str:
    # Prefer S3 if configured; else save to local artifacts dir
    if settings.S3_BUCKET and aioboto3:
        session = aioboto3.Session()
        async with session.client("s3", region_name=settings.AWS_REGION) as s3:
            await s3.put_object(Bucket=settings.S3_BUCKET, Key=key, Body=data, ContentType=content_type)
            url = await s3.generate_presigned_url(
                "get_object",
                Params={"Bucket": settings.S3_BUCKET, "Key": key},
                ExpiresIn=3600,
            )
            return url

    # ---- local fallback ----
    base_dir = settings.ARTIFACTS_DIR.rstrip("/")
    # allow callers to pass either "artifacts/x/y.pptx" or "x/y.pptx"
    rel = key.split("artifacts/", 1)[1] if key.startswith("artifacts/") else key
    path = os.path.join(base_dir, rel)
    os.makedirs(os.path.dirname(path), exist_ok=True)
    with open(path, "wb") as f:
        f.write(data)

    # Build URL (absolute if PUBLIC_BASE_URL is set)
    base_url = (settings.PUBLIC_BASE_URL or "").rstrip("/")
    url_path = f"/artifacts/{rel}"
    return f"{base_url}{url_path}" if base_url else url_path


# ---------- tools (called by the LLM agent) ----------

async def load_pdf_text(
    source_url: Optional[str] = None,
    s3_key: Optional[str] = None,
    max_chars: int = 40_000,
    *,
    tenant_id: Optional[str] = None,
    job_id: Optional[str] = None,
) -> dict:
    """
    Fetch a PDF from a public URL or S3 and return extracted text for outlining.

    Args:
      source_url: HTTPS URL to the PDF.
      s3_key: S3 key within your configured bucket.
      max_chars: Truncate extracted text to this many characters.
      tenant_id: (injected) Tenant context; not used directly by the tool.
      job_id: (injected) Job context; not used directly by the tool.

    Returns:
      {"text": "...", "truncated": bool, "chars": int}
    """
    if not (source_url or s3_key):
        return {"text": "", "truncated": False, "chars": 0}

    data = await (
        _fetch_bytes_from_url(source_url) if source_url else _fetch_bytes_from_s3(s3_key)  # type: ignore[arg-type]
    )

    text = await _extract_text_async(data)
    truncated = False
    if len(text) > max_chars:
        text = text[:max_chars]
        truncated = True

    return {"text": text, "truncated": truncated, "chars": len(text)}


def _add_title_slide(prs: Presentation, title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title
    slide.shapes.title.text = title or "Presentation"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = ""


def _add_bullet_slide(prs: Presentation, title: str, bullets: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide.shapes.title.text = (title or "")[:120]
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for b in bullets or []:
        p = tf.add_paragraph()
        p.text = b[:200]
        p.font.size = Pt(18)
        p.level = 0


async def build_pptx(
    outline_json: str,
    title: Optional[str] = None,
    theme: str = "light",
    output_key: Optional[str] = None,
    *,
    tenant_id: Optional[str] = None,
    job_id: Optional[str] = None,
) -> dict:
    """
    Build a PPTX from an outline (JSON string). Upload to S3 (if configured) or save locally.

    Args:
      outline_json: JSON array like [{"title":"Intro","bullets":["...","..."]}, ...]
      title: Optional deck title.
      theme: Reserved for future theming; currently no-op ("light"|"dark").
      output_key: Explicit S3/local key, otherwise artifacts/{tenant}/{job}.pptx.
      tenant_id: (injected) Tenant context; used to compute default key.
      job_id: (injected) Job context; used to compute default key.

    Returns:
      {"pptx_s3_key": str, "pptx_url": str, "slides": int, "outline": list}
    """
    try:
        outline = json.loads(outline_json)
        if not isinstance(outline, list):
            raise ValueError("outline_json must decode to a list")
    except Exception as e:
        raise ValueError(f"Invalid outline_json: {e}") from e

    prs = Presentation()
    _add_title_slide(prs, title or ((outline[0].get("title") or "Presentation") if outline else "Presentation"))

    for item in outline:
        item_title = (item.get("title") or "").strip() if isinstance(item, dict) else ""
        item_bullets = item.get("bullets") if isinstance(item, dict) else []
        if not isinstance(item_bullets, list):
            item_bullets = []
        _add_bullet_slide(prs, item_title, item_bullets)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    key = output_key or f"{tenant_id or 'demo'}/{job_id or uuid.uuid4()}.pptx"
    url = await _upload_bytes(buf.getvalue(), key)

    return {
        "pptx_s3_key": key,
        "pptx_url": url,
        "slides": len(prs.slides),
        "outline": outline,
    }
